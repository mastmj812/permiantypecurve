"""Deal dossier PPTX — a working discussion deck, not a final exhibit.

Structure: one slide per narvi scenario (plan-view map left, gunbarrel
right, well-count subtitle), then the deal's type curves rendered
exactly like the existing per-curve slide export (param table + rate /
cum charts + cohort map, one slide per stream). The wells-table slide
and probit are deliberately omitted — this deck is for talking through
a deal while work is ongoing.

All raster panels arrive from the client (the dossier preview page
captures its SVG charts and MapLibre canvases), mirroring the
type-curve slide export: the backend only assembles the deck from the
brand template, so charts and maps stay pixel-identical to the app.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Any
from uuid import UUID

from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.db.models import TypeCurve
from app.exports.blueox import RATIO_REFUSAL_NOTE, ratio_mode_streams
from app.exports.pptx_builder import (
    _STREAM_TITLE,
    TEMPLATE_PATH,
    _duplicate_slide,
    _fill_param_table,
    _find_table,
    _place_chart_images,
    _set_title_text,
)

# Scenario-slide geometry (16:9, 13.333" x 7.5"): two side-by-side
# panels under the title, mirroring the chart column's 0.64" margins.
_MARGIN_IN = 0.64
_SLIDE_WIDTH_IN = 13.333
_PANEL_TOP_IN = 1.75
_PANEL_HEIGHT_IN = 4.9
_PANEL_GAP_IN = 0.15
_PANEL_WIDTH_IN = (_SLIDE_WIDTH_IN - 2 * _MARGIN_IN - _PANEL_GAP_IN) / 2
_SUBTITLE_TOP_IN = 1.32
_SUBTITLE_FONT_PT = 11

# Matching pixel dimensions for the client's capture panels (96 px/in)
# so the PNG aspect equals the picture box and nothing stretches.
SCENARIO_PANEL_PX = (
    round(_PANEL_WIDTH_IN * 96),
    round(_PANEL_HEIGHT_IN * 96),
)

# TC-vs-Novi comparison slide: one full-width figure (the 6-panel
# rate/cum grid) under the title + subtitle.
_COMPARISON_TOP_IN = 1.6
_COMPARISON_HEIGHT_IN = 5.5
_COMPARISON_WIDTH_IN = _SLIDE_WIDTH_IN - 2 * _MARGIN_IN
COMPARISON_PANEL_PX = (
    round(_COMPARISON_WIDTH_IN * 96),
    round(_COMPARISON_HEIGHT_IN * 96),
)


@dataclass(frozen=True)
class ScenarioSlideInput:
    """One narvi scenario's slide: client-captured map + gunbarrel."""

    title: str
    subtitle: str  # bench/category well counts, built client-side
    map_png: bytes
    gunbarrel_png: bytes


@dataclass(frozen=True)
class ComparisonSlideInput:
    """One zone's TC-vs-Novi comparison slide: the client-captured
    6-panel figure (oil/gas/water x rate/cum, TC band vs Novi median)."""

    title: str
    subtitle: str  # alignment/normalization/risking basis, built client-side
    figure_png: bytes


@dataclass(frozen=True)
class CurveSlideInput:
    """One type curve's stream slides (same panels as the slide export)."""

    type_curve_id: UUID
    # stream -> (rate_png, cum_png); must carry oil, gas, water.
    stream_pngs: dict[str, tuple[bytes, bytes]] = field(default_factory=dict)
    map_png: bytes = b""


def build_deal_dossier_pptx(
    session: Session,
    scenarios: list[ScenarioSlideInput],
    curves: list[CurveSlideInput],
    comparisons: list[ComparisonSlideInput] | None = None,
) -> bytes:
    """Assemble the dossier deck from the brand template.

    Slide order: scenarios (one each), then per curve the oil / gas /
    water stream slides, then one TC-vs-Novi comparison slide per zone
    that has one. Raises ValueError on an unknown type curve or a
    missing stream panel.
    """
    if not scenarios and not curves:
        raise ValueError("dossier needs at least one scenario or curve")

    pres = Presentation(str(TEMPLATE_PATH))
    if len(pres.slides) < 2:
        raise ValueError("template must have at least 2 slides (stream slide + well table)")

    # Template is [stream_template, wells]. Every dossier slide starts
    # as a duplicate of the stream template (appended at the end); the
    # two template slides are deleted once all content slides exist.
    for sc in scenarios:
        _duplicate_slide(pres, source_idx=0)
        _build_scenario_slide(pres.slides[-1], sc)

    streams = ["oil", "gas", "water"]
    for cv in curves:
        tc = session.get(TypeCurve, cv.type_curve_id)
        if tc is None:
            raise ValueError(f"type curve {cv.type_curve_id} not found")
        # The dossier's param table declares Arps params per stream — a
        # ratio-mode stream has none. HARD refusal, same posture as the
        # Blue Ox drop builder (in-app use stays unrestricted).
        ratio_streams = ratio_mode_streams(tc.series)
        if ratio_streams:
            raise ValueError(
                f"curve {tc.name}: stream(s) {', '.join(ratio_streams)} are "
                f"ratio-mode — {RATIO_REFUSAL_NOTE}"
            )
        for stream in streams:
            if stream not in cv.stream_pngs:
                raise ValueError(f"curve {tc.name}: missing {stream} panels")
            _duplicate_slide(pres, source_idx=0)
            slide = pres.slides[-1]
            rate_png, cum_png = cv.stream_pngs[stream]
            _set_title_text(slide, f"{tc.name} {_STREAM_TITLE[stream]}")
            _fill_param_table(_find_table(slide), tc, None)
            _place_chart_images(slide, rate_png, cum_png, cv.map_png)

    for cp in comparisons or []:
        _duplicate_slide(pres, source_idx=0)
        _build_comparison_slide(pres.slides[-1], cp)

    # Drop the template slides (indices 0 and 1) now that the content
    # slides are in place — delete back-to-front so indices hold.
    _delete_slide(pres, 1)
    _delete_slide(pres, 0)

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


def _build_scenario_slide(slide: Slide, sc: ScenarioSlideInput) -> None:
    """Turn a duplicated stream-template slide into a scenario slide:
    keep the title, strip the param table and template picture, place
    map (left) + gunbarrel (right) with a subtitle line between."""
    _set_title_text(slide, sc.title)
    for shape in list(slide.shapes):
        is_table = isinstance(shape, GraphicFrame) and shape.has_table
        if is_table or isinstance(shape, Picture):
            shape._element.getparent().remove(shape._element)

    if sc.subtitle:
        tb = slide.shapes.add_textbox(
            Inches(_MARGIN_IN),
            Inches(_SUBTITLE_TOP_IN),
            Inches(_SLIDE_WIDTH_IN - 2 * _MARGIN_IN),
            Inches(0.3),
        )
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = sc.subtitle
        run.font.size = Pt(_SUBTITLE_FONT_PT)

    slide.shapes.add_picture(
        io.BytesIO(sc.map_png),
        Inches(_MARGIN_IN),
        Inches(_PANEL_TOP_IN),
        width=Inches(_PANEL_WIDTH_IN),
        height=Inches(_PANEL_HEIGHT_IN),
    )
    slide.shapes.add_picture(
        io.BytesIO(sc.gunbarrel_png),
        Inches(_MARGIN_IN + _PANEL_WIDTH_IN + _PANEL_GAP_IN),
        Inches(_PANEL_TOP_IN),
        width=Inches(_PANEL_WIDTH_IN),
        height=Inches(_PANEL_HEIGHT_IN),
    )


def _build_comparison_slide(slide: Slide, cp: ComparisonSlideInput) -> None:
    """Title + subtitle + one full-width 6-panel comparison figure."""
    _set_title_text(slide, cp.title)
    for shape in list(slide.shapes):
        is_table = isinstance(shape, GraphicFrame) and shape.has_table
        if is_table or isinstance(shape, Picture):
            shape._element.getparent().remove(shape._element)

    if cp.subtitle:
        tb = slide.shapes.add_textbox(
            Inches(_MARGIN_IN),
            Inches(_SUBTITLE_TOP_IN),
            Inches(_SLIDE_WIDTH_IN - 2 * _MARGIN_IN),
            Inches(0.3),
        )
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = cp.subtitle
        run.font.size = Pt(_SUBTITLE_FONT_PT)

    slide.shapes.add_picture(
        io.BytesIO(cp.figure_png),
        Inches(_MARGIN_IN),
        Inches(_COMPARISON_TOP_IN),
        width=Inches(_COMPARISON_WIDTH_IN),
        height=Inches(_COMPARISON_HEIGHT_IN),
    )


def _delete_slide(pres: Any, idx: int) -> None:
    """Remove a slide from the deck. python-pptx has no delete API; the
    documented workaround is to drop the slide's relationship and its
    entry in the slide-id list."""
    slide = pres.slides[idx]
    # Find the rId for this slide part and drop it.
    for rel_id, rel in list(pres.part.rels.items()):
        if rel.target_part is slide.part:
            pres.part.drop_rel(rel_id)
            break
    sld_id_lst = pres.slides._sldIdLst
    sld_id_lst.remove(list(sld_id_lst)[idx])
