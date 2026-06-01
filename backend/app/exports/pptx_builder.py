"""PowerPoint export for a saved type curve + optional comparison curve.

Loads the user's brand template at ``templates/deal_slide.pptx`` and
populates it in place:

  * Slide 1 — title text, 17-column param table, single composite
    picture (rate + cum + map, snapshotted by the client from the
    in-app slide preview).
  * Slide 2 — 12-column well table, one row per included api10.

Shape discovery is by **type**, not by name, because the template
uses PowerPoint's auto-generated shape names ("Title 1", "Table 6",
etc.) which aren't stable across edits. As long as each slide has
the expected single placeholder / table / picture, this keeps
working through template tweaks.

Cell-text rewrites use the find-first-run approach so paragraph
alignment and run-level formatting (font size, bold, lang) carry
forward. Row add/remove on tables manipulates ``<a:tr>`` via
``copy.deepcopy`` so new rows inherit the formatting of an existing
data row — a pattern documented in the python-pptx docs as the
canonical way to grow / shrink a table since the library has no
high-level row API.
"""

from __future__ import annotations

import copy
import io
from pathlib import Path
from typing import Any
from uuid import UUID

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.table import Table, _Cell
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session


# Default font size we fall back to when a cell's first paragraph has
# no runs (sample template rows occasionally have trailing cells with
# empty paragraphs, e.g. wells missing forecast data). Without this,
# python-pptx's add_run() creates a run with no rPr → PowerPoint falls
# back to its 18pt default, which sticks out next to the table's body
# rows.
_DEFAULT_BODY_FONT_PT = 8

from app.db.models import TypeCurve
from app.exports.param_row import PARAM_HEADERS, format_param_row
from app.exports.well_rows import PER_WELL_HEADERS, per_well_rows


TEMPLATE_PATH = Path(__file__).parent / "templates" / "deal_slide.pptx"


# Slide-level positions for chart pictures. Each chart picture is
# EXACTLY 5.42" × 2.16" on the slide — to guarantee that survives the
# export we delete the template's original single picture frame and
# place fresh pictures at these absolute coordinates.
#
# Two layouts:
#   Without probit (default): Rate top-left, Cum bottom-left, Map
#     full-height right (6.93" × 4.36" — wider, fills the right
#     column).
#   With probit: Rate top-left, Cum bottom-left, Map top-right, Probit
#     bottom-right. Map shrinks to 5.42" × 2.16" to match the chart
#     panel size; probit takes the remaining bottom-right cell.
_CHART_WIDTH_IN = 5.42
_CHART_HEIGHT_IN = 2.16
_CHART_LEFT_IN = 0.64
_CHART_TOP_IN = 2.28
_CHART_GAP_IN = 0.04
_RIGHT_COL_LEFT_IN = _CHART_LEFT_IN + _CHART_WIDTH_IN + 0.08
# Big map (probit excluded): fills the right column entirely.
_BIG_MAP_WIDTH_IN = 6.93
_BIG_MAP_HEIGHT_IN = _CHART_HEIGHT_IN * 2 + _CHART_GAP_IN

# Footnote (downtime-filter disclosure) sits below the cum chart,
# left-aligned with the chart's y-axis. The slide chart uses
# TypeCurveChart with compactLayout=true → PAD.left = 50 px on a
# 520 px-wide SVG, so the y-axis sits at 50/520 of the chart's
# 5.42" width = ~0.52" from the chart panel's left edge.
_FOOTNOTE_TOP_IN = _CHART_TOP_IN + 2 * _CHART_HEIGHT_IN + _CHART_GAP_IN
_FOOTNOTE_LEFT_IN = _CHART_LEFT_IN + 50.0 / 520.0 * _CHART_WIDTH_IN
_FOOTNOTE_WIDTH_IN = _CHART_WIDTH_IN  # generous; PowerPoint won't wrap unless overflowing
_FOOTNOTE_HEIGHT_IN = 0.18
_FOOTNOTE_FONT_PT = 7
_FOOTNOTE_TEXT = "Actual production filtered to non-downtime months."


_STREAM_TITLE: dict[str, str] = {
    "oil": "Oil",
    "gas": "Gas",
    "water": "Water",
}


def build_deal_slide_pptx(
    session: Session,
    tc_id: UUID,
    compare_id: UUID | None,
    *,
    stream_pngs: dict[str, tuple[bytes, bytes]],
    map_png: bytes,
    probit_pngs: dict[str, bytes] | None = None,
) -> bytes:
    """Build a 4-slide .pptx: Oil, Gas, Water, Type Curve Wells.

    ``stream_pngs`` maps each stream name → ``(rate_png, cum_png)``.
    ``probit_pngs`` (optional) maps each stream name → probit_png; when
    provided each stream slide gets a probit chart bottom-right and the
    map shrinks to match the chart panel size. When None, the map fills
    the right column as before.
    """
    tc = session.get(TypeCurve, tc_id)
    if tc is None:
        raise ValueError(f"type curve {tc_id} not found")
    prev = session.get(TypeCurve, compare_id) if compare_id is not None else None

    pres = Presentation(str(TEMPLATE_PATH))
    if len(pres.slides) < 2:
        raise ValueError(
            "template must have at least 2 slides (stream slide + well table)"
        )

    # Template starts as [stream_template (slide 0), wells (slide 1)].
    # Duplicate the stream slide twice — both copies land at the end —
    # then reorder to [oil, gas, water, wells].
    _duplicate_slide(pres, source_idx=0)  # → [oil, wells, gas_copy]
    _duplicate_slide(pres, source_idx=0)  # → [oil, wells, gas_copy, water_copy]
    # Move the wells slide (currently at index 1) to the end.
    _move_slide(pres, from_idx=1, to_idx=3)
    # Slides are now [oil, gas, water, wells].

    streams = ["oil", "gas", "water"]
    for slide_idx, stream in enumerate(streams):
        slide = pres.slides[slide_idx]
        rate_png, cum_png = stream_pngs[stream]
        _set_title_text(slide, f"{tc.name} {_STREAM_TITLE[stream]}")
        # The param table is the SAME 17-col layout on all three
        # stream slides (it already shows oil + gas + water columns
        # side-by-side; the slide-level distinction is the charts).
        param_table = _find_table(slide)
        _fill_param_table(param_table, tc, prev)
        probit_png = (probit_pngs or {}).get(stream)
        _place_chart_images(
            slide, rate_png, cum_png, map_png, probit_png=probit_png,
        )

    # --- wells slide (now at index 3) ---
    wells_slide = pres.slides[3]
    well_table = _find_table(wells_slide)
    rows = per_well_rows(session, list(tc.included_api10s or []))
    _fill_well_table(well_table, rows)

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()


# ============================ slide duplication ============================


def _duplicate_slide(pres: Any, source_idx: int) -> None:
    """Append a deep copy of ``pres.slides[source_idx]`` to the end of
    the deck. python-pptx has no built-in clone-slide API; the
    canonical workaround is to add a new slide using the source's
    layout, strip the placeholders the layout drops in, then deep-
    copy each shape from the source into the new slide's spTree.

    Picture relationships are re-linked to the new slide's part so
    image data still resolves after the copy.
    """
    from copy import deepcopy

    source = pres.slides[source_idx]
    layout = source.slide_layout
    dest = pres.slides.add_slide(layout)

    # Remove placeholder shapes the layout dropped in — we want a clean
    # canvas to deep-copy the source shapes onto.
    for ph in list(dest.shapes):
        ph.element.getparent().remove(ph.element)

    spTree = dest.shapes._spTree
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        spTree.append(new_el)

    # Re-link any picture / chart / OLE relationships from source to
    # dest. The deep-copied XML carries rId references that point to
    # the SOURCE slide's part; without re-linking, openpyxl/PowerPoint
    # can't resolve them. For our template the only relationship-bound
    # shape is the picture (rId2 → ppt/media/image1.png), which the
    # builder removes immediately anyway (replaced with the user's
    # chart PNGs). But re-linking is robust if the template ever grows
    # a logo or chart.
    for rel in source.part.rels.values():
        if rel.reltype not in dest.part.rels:
            try:
                dest.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                # Best-effort — a missing relationship at this stage
                # only matters if the cloned slide depends on it.
                pass


def _move_slide(pres: Any, from_idx: int, to_idx: int) -> None:
    """Move a slide so its post-move index is ``to_idx``. python-pptx
    exposes the slide-id list directly; reorder via lxml insert.

    Convention: ``to_idx`` is the FINAL post-move position in the deck
    (0-indexed). lxml ``insert(idx, ...)`` clamps idx ≥ len to "append",
    so moving to the end is just ``to_idx = len(slides) - 1``.
    """
    sld_id_lst = pres.slides._sldIdLst
    slides = list(sld_id_lst)
    moving = slides[from_idx]
    sld_id_lst.remove(moving)
    sld_id_lst.insert(to_idx, moving)


# ============================ shape discovery ============================


def _find_title_shape(slide: Slide) -> Shape | None:
    # The title is the single placeholder with idx=0 (or PP_PLACEHOLDER.TITLE).
    # python-pptx exposes `slide.shapes.title` directly when the slide
    # layout has one — but our template may not declare it that way.
    if slide.shapes.title is not None:
        return slide.shapes.title
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            return shape
    return None


def _find_table(slide: Slide) -> Table:
    for shape in slide.shapes:
        if isinstance(shape, GraphicFrame) and shape.has_table:
            return shape.table
    raise ValueError(f"no table on slide {slide.slide_id}")


# ============================ title ============================


def _set_title_text(slide: Slide, text: str) -> None:
    title = _find_title_shape(slide)
    if title is None:
        return
    _set_text_frame_text(title.text_frame, text)


def _set_text_frame_text(tf: Any, text: str) -> None:
    """Replace text in a text frame without nuking paragraph / run
    formatting. Targets the first paragraph's first run; other runs in
    that paragraph are removed, other paragraphs are left untouched.

    When the first paragraph has NO existing runs (template cell was
    empty), the new run we add carries no rPr → PowerPoint renders it
    at its 18pt default which towers over the surrounding 8pt body
    text. Pin the new run to ``_DEFAULT_BODY_FONT_PT`` so it visually
    matches the rest of the table.
    """
    text = "" if text is None else str(text)
    paragraphs = tf.paragraphs
    if not paragraphs:
        tf.text = text
        return
    p = paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for extra in list(p.runs[1:]):
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
        run.font.size = Pt(_DEFAULT_BODY_FONT_PT)
        run.text = text


# ============================ tables ============================


def _set_cell_text(cell: _Cell, text: str) -> None:
    _set_text_frame_text(cell.text_frame, text)


def _write_row(table: Table, row_idx: int, values: tuple[Any, ...]) -> None:
    """Write `values` into the existing row's cells (preserves formatting)."""
    row = table.rows[row_idx]
    for col_idx, val in enumerate(values):
        if col_idx >= len(row.cells):
            break
        cell = row.cells[col_idx]
        _set_cell_text(cell, _format_cell(val))


def _format_cell(val: Any) -> str:
    if val is None:
        return ""
    return str(val)


def _clone_row(table: Table, source_row_idx: int) -> None:
    """Append a new row to the table by deep-copying an existing row.

    The clone inherits the source row's cell formatting (font size,
    bold, alignment, fills). python-pptx has no high-level row-add
    API; this is the documented workaround.
    """
    src_tr = table.rows[source_row_idx]._tr
    new_tr = copy.deepcopy(src_tr)
    table._tbl.append(new_tr)


def _delete_row(table: Table, row_idx: int) -> None:
    tr = table.rows[row_idx]._tr
    tr.getparent().remove(tr)


def _fill_param_table(
    table: Table, tc: TypeCurve, prev: TypeCurve | None
) -> None:
    """Slide-1 param table: row 0 is header, row 1 is the current
    curve's row, optional row 2 is the previous curve when comparing.

    If the template has more or fewer data rows than what we need,
    delete or clone to converge — the template ships with exactly 2
    rows (header + 1 data) so most exports will hit the "modify in
    place" or "modify + clone one row" path.
    """
    header_rows = 1  # row 0 is the header
    needed_data_rows = 2 if prev is not None else 1
    current_data_rows = len(table.rows) - header_rows

    # Sanity-check column count matches; if the template changes shape
    # we'd silently mismatch, so be loud.
    expected_cols = len(PARAM_HEADERS)
    actual_cols = len(table.columns)
    if actual_cols != expected_cols:
        raise ValueError(
            f"param table column mismatch: template has {actual_cols}, "
            f"formatters produce {expected_cols}"
        )

    # Grow or shrink the data-row count by deep-copying / removing.
    while current_data_rows < needed_data_rows:
        _clone_row(table, source_row_idx=header_rows)  # clone the existing data row
        current_data_rows += 1
    while current_data_rows > needed_data_rows:
        _delete_row(table, row_idx=header_rows + current_data_rows - 1)
        current_data_rows -= 1

    _write_row(table, row_idx=header_rows, values=format_param_row(tc))
    if prev is not None:
        _write_row(
            table, row_idx=header_rows + 1, values=format_param_row(prev)
        )


def _fill_well_table(table: Table, rows: list[tuple[Any, ...]]) -> None:
    """Slide-2 well table: row 0 is header, rows 1..N are one well each.

    Add or remove rows to match the cohort size, then write each cell
    via the formatter map (Oil/Gas/Oil EUR/ft as comma-thousands int,
    GOR/WOR as 1-decimal float). Mirrors ``PER_WELL_COL_FORMATS`` from
    the xlsx export so PPT and xlsx render identically.
    """
    header_rows = 1
    needed = len(rows)
    current_data_rows = len(table.rows) - header_rows

    # The template ships with sample data; we need one row per cohort
    # well. Grow/shrink via the clone/delete helpers.
    while current_data_rows < needed:
        _clone_row(table, source_row_idx=header_rows)
        current_data_rows += 1
    while current_data_rows > needed:
        _delete_row(table, row_idx=header_rows + current_data_rows - 1)
        current_data_rows -= 1

    for i, row in enumerate(rows):
        formatted = tuple(_format_well_cell(i_col, v) for i_col, v in enumerate(row))
        _write_row(table, row_idx=header_rows + i, values=formatted)


def _format_well_cell(col_idx: int, val: Any) -> str:
    """Per-column number formatting for the well table.

    Mirrors ``PER_WELL_COL_FORMATS`` in ``app.exports.well_rows``:
    columns 5–9 (BWPF, PPF, Oil EUR, Gas EUR, Oil EUR/ft) display as
    comma-thousands integers; columns 10–11 (GOR EUR, WOR EUR) as
    one decimal; everything else as raw text.
    """
    if val is None:
        return ""
    # 0-indexed: api10=0, Wellname=1, Operator=2, Formation=3, Lateral=4,
    # BWPF=5, PPF=6, Oil EUR=7, Gas EUR=8, Oil EUR/ft=9, GOR=10, WOR=11
    if col_idx == 4:  # Lateral Length — display as plain integer
        try:
            return f"{round(float(val)):,}"
        except (TypeError, ValueError):
            return str(val)
    if col_idx in (5, 6, 7, 8, 9):
        try:
            return f"{round(float(val)):,}"
        except (TypeError, ValueError):
            return str(val)
    if col_idx in (10, 11):
        try:
            return f"{float(val):.1f}"
        except (TypeError, ValueError):
            return str(val)
    return str(val)


# ============================ pictures ============================


def _add_downtime_footnote(slide: Slide) -> None:
    """Place the small downtime-disclosure footnote below the cum
    chart, left-aligned with the chart y-axis. 7pt text, single line.
    Doesn't reflow or resize any existing shape — additive only.
    """
    tb = slide.shapes.add_textbox(
        Inches(_FOOTNOTE_LEFT_IN),
        Inches(_FOOTNOTE_TOP_IN),
        Inches(_FOOTNOTE_WIDTH_IN),
        Inches(_FOOTNOTE_HEIGHT_IN),
    )
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _FOOTNOTE_TEXT
    run.font.size = Pt(_FOOTNOTE_FONT_PT)


def _place_chart_images(
    slide: Slide,
    rate_png: bytes,
    cum_png: bytes,
    map_png: bytes,
    *,
    probit_png: bytes | None = None,
) -> None:
    """Place rate/cum/map (+ optional probit) on a stream slide.

    Without probit: rate top-left, cum bottom-left, map fills right
    column (6.93" × 4.36"). With probit: rate TL, cum BL, map TR
    (5.42" × 2.16"), probit BR (same). The template's original single
    picture frame is removed first either way; new picture shapes go
    in its place at exact inch positions so chart sizes are
    guaranteed regardless of any template frame's aspect.
    """
    # Find and remove the template's original picture (if any). If
    # the user later updates the template to drop the placeholder
    # entirely, this is a no-op.
    for shape in list(slide.shapes):
        if isinstance(shape, Picture):
            shape._element.getparent().remove(shape._element)

    chart_left = Inches(_CHART_LEFT_IN)
    chart_width = Inches(_CHART_WIDTH_IN)
    chart_height = Inches(_CHART_HEIGHT_IN)
    rate_top = Inches(_CHART_TOP_IN)
    cum_top = Inches(_CHART_TOP_IN + _CHART_HEIGHT_IN + _CHART_GAP_IN)

    slide.shapes.add_picture(
        io.BytesIO(rate_png), chart_left, rate_top,
        width=chart_width, height=chart_height,
    )
    slide.shapes.add_picture(
        io.BytesIO(cum_png), chart_left, cum_top,
        width=chart_width, height=chart_height,
    )
    _add_downtime_footnote(slide)

    right_col_left = Inches(_RIGHT_COL_LEFT_IN)

    if probit_png is None:
        # Big-map layout — map fills the right column.
        slide.shapes.add_picture(
            io.BytesIO(map_png),
            right_col_left, Inches(_CHART_TOP_IN),
            width=Inches(_BIG_MAP_WIDTH_IN), height=Inches(_BIG_MAP_HEIGHT_IN),
        )
        return

    # With-probit layout — map shrinks to the chart panel size, probit
    # goes bottom-right at the same size.
    slide.shapes.add_picture(
        io.BytesIO(map_png),
        right_col_left, rate_top,
        width=chart_width, height=chart_height,
    )
    slide.shapes.add_picture(
        io.BytesIO(probit_png),
        right_col_left, cum_top,
        width=chart_width, height=chart_height,
    )
