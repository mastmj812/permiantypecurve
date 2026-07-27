"""Deck-shape tests for the deal dossier export.

Builds the deck bytes directly via ``build_deal_dossier_pptx`` with a
stub session (no DB): scenario slides first (map + gunbarrel picture
pair, no param table), then oil/gas/water stream slides per curve in
the same layout as the type-curve slide export.
"""

from __future__ import annotations

import io
import uuid
from typing import Any

import pytest
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture

from app.exports.dossier import (
    CurveSlideInput,
    ScenarioSlideInput,
    build_deal_dossier_pptx,
)
from tests.test_deal_export import _curve

# Minimal valid 1x1 PNG — python-pptx parses image headers on
# add_picture, so the bytes must be a real image.
_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01"
    b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


class _StubSession:
    """Just enough of a Session for the builder's ``session.get``."""

    def __init__(self, curves: list[Any]) -> None:
        self._by_id = {c.id: c for c in curves}

    def get(self, model: Any, key: Any) -> Any:
        return self._by_id.get(key)


def _scenario(title: str) -> ScenarioSlideInput:
    return ScenarioSlideInput(
        title=title,
        subtitle="PUD WCB_2 x2 - PDP x4",
        map_png=_PNG,
        gunbarrel_png=_PNG,
    )


def _curve_input(tc: Any) -> CurveSlideInput:
    return CurveSlideInput(
        type_curve_id=tc.id,
        stream_pngs=dict.fromkeys(("oil", "gas", "water"), (_PNG, _PNG)),
        map_png=_PNG,
    )


def _slide_texts(slide: Any) -> str:
    return " ".join(
        sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    )


def test_dossier_slide_order_and_shapes() -> None:
    tc = _curve("holdTheLine_wca_v1")
    content = build_deal_dossier_pptx(
        _StubSession([tc]),  # type: ignore[arg-type]
        [_scenario("plan_brotime_20"), _scenario("plan_brotime_20_35")],
        [_curve_input(tc)],
    )
    pres = Presentation(io.BytesIO(content))
    # 2 scenario slides + 3 stream slides; template slides removed.
    assert len(pres.slides) == 5

    for idx, title in ((0, "plan_brotime_20"), (1, "plan_brotime_20_35")):
        slide = pres.slides[idx]
        pictures = [s for s in slide.shapes if isinstance(s, Picture)]
        tables = [
            s for s in slide.shapes
            if isinstance(s, GraphicFrame) and s.has_table
        ]
        assert len(pictures) == 2  # map + gunbarrel
        assert not tables  # param table stripped
        text = _slide_texts(slide)
        assert title in text
        assert "PUD WCB_2 x2" in text  # subtitle carried through

    for idx, stream in ((2, "Oil"), (3, "Gas"), (4, "Water")):
        slide = pres.slides[idx]
        tables = [
            s for s in slide.shapes
            if isinstance(s, GraphicFrame) and s.has_table
        ]
        assert len(tables) == 1  # 17-col param table
        pictures = [s for s in slide.shapes if isinstance(s, Picture)]
        assert len(pictures) == 3  # rate + cum + map
        assert f"holdTheLine_wca_v1 {stream}" in _slide_texts(slide)


def test_dossier_scenarios_only() -> None:
    content = build_deal_dossier_pptx(
        _StubSession([]),  # type: ignore[arg-type]
        [_scenario("plan_castaway_23")],
        [],
    )
    pres = Presentation(io.BytesIO(content))
    assert len(pres.slides) == 1


def test_dossier_rejects_unknown_curve_and_empty() -> None:
    with pytest.raises(ValueError, match="not found"):
        build_deal_dossier_pptx(
            _StubSession([]),  # type: ignore[arg-type]
            [],
            [
                CurveSlideInput(
                    type_curve_id=uuid.uuid4(),
                    stream_pngs=dict.fromkeys(
                        ("oil", "gas", "water"), (_PNG, _PNG)
                    ),
                    map_png=_PNG,
                )
            ],
        )
    with pytest.raises(ValueError, match="at least one"):
        build_deal_dossier_pptx(_StubSession([]), [], [])  # type: ignore[arg-type]


def test_dossier_comparison_slides_append_after_curves() -> None:
    from app.exports.dossier import ComparisonSlideInput

    tc = _curve("holdTheLine_wca_v1")
    content = build_deal_dossier_pptx(
        _StubSession([tc]),  # type: ignore[arg-type]
        [_scenario("plan_brotime_20")],
        [_curve_input(tc)],
        [
            ComparisonSlideInput(
                title="WCA — Type Curve vs Novi ML (n=4: 4 PUD / 0 RES)",
                subtitle="per 1,000 ft lateral · TC aligned to peak, Novi to IP",
                figure_png=_PNG,
            ),
        ],
    )
    pres = Presentation(io.BytesIO(content))
    # 1 scenario + 3 stream slides + 1 comparison slide.
    assert len(pres.slides) == 5
    slide = pres.slides[4]
    pictures = [s for s in slide.shapes if isinstance(s, Picture)]
    tables = [
        s for s in slide.shapes
        if isinstance(s, GraphicFrame) and s.has_table
    ]
    assert len(pictures) == 1  # the single full-width figure
    assert not tables  # param table stripped
    text = _slide_texts(slide)
    assert "Type Curve vs Novi ML" in text
    assert "Novi to IP" in text
